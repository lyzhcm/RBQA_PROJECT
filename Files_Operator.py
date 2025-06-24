import os
import json
import pandas as pd
import PyPDF2
import tempfile
import pptx
from docx import Document
import streamlit as st
import hashlib
import re
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from sentence_transformers import SentenceTransformer

def parse_file(file):
    """解析上传的文件内容"""
    content = ""
    file_type = file.name.split(".")[-1].lower()

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
            tmp.write(file.getvalue())
            tmp_path = tmp.name

        if file_type == "txt":
            with open(tmp_path, "r", encoding="utf-8") as f:
                content = f.read()

        elif file_type == "pdf":
            reader = PyPDF2.PdfReader(tmp_path)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    content += text + "\n"

        elif file_type == "docx":
            doc = Document(tmp_path)
            for para in doc.paragraphs:
                content += para.text + "\n"

        elif file_type == "pptx":
            prs = pptx.Presentation(tmp_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"

        os.unlink(tmp_path)
        return content

    except Exception as e:
        st.error(f"解析文件时出错: {str(e)}")
        return None
    
# 文本预处理函数
def preprocess_text(text):
    # 转换为小写
    text = text.lower()
    # 移除特殊字符和数字
    text = re.sub(r'[^a-zA-Z\u4e00-\u9fff\s]', '', text)
    # 分词
    tokens = word_tokenize(text)
    # 移除停用词
    stop_words = set(stopwords.words('english') + stopwords.words('chinese'))
    tokens = [word for word in tokens if word not in stop_words]
    return " ".join(tokens)

# 将文档内容分割为适合处理的片段
def split_into_chunks(text, chunk_size=200):
    chunks = []
    words = text.split()

    for i in range(0, len(words), chunk_size):
        chunks.append(" ".join(words[i:i + chunk_size]))

    return chunks


# 生成文件唯一ID
def generate_file_id(file_content):
    return hashlib.md5(file_content).hexdigest()



# 加载嵌入模型（缓存避免重复加载）
@st.cache_resource
def load_embedding_model():
    return SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')


# 示例用法
# reader = Reader()
# content = reader.read('example.txt')
# print(content)