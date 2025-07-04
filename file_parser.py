#file_parser
import os
import json
import pandas as pd
import PyPDF2
import tempfile
from pptx import Presentation  # 修复导入方式
from docx import Document
import streamlit as st
import hashlib
import re
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from file_registry import FileRegistry
from config import PERSISTENT_UPLOAD_FOLDER
from pathlib import Path
from werkzeug.utils import secure_filename
import nltk

# 下载nltk资源
try:
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('punkt')
    nltk.download('stopwords')

def save_uploaded_file(file, file_id):
    """持久化保存文件"""
    try:
        # 确保文件名安全且唯一
        safe_name = f"{file_id}_{secure_filename(file.name)}"
        save_path = os.path.join(PERSISTENT_UPLOAD_FOLDER, safe_name)
        
        # 保存文件
        with open(save_path, "wb") as f:
            f.write(file.getbuffer())
            
        # 注册文件
        FileRegistry.add_file(file_id, file.name, save_path)
        return save_path
    except Exception as e:
        st.error(f"文件保存失败: {str(e)}")
        return None
    
def parse_file(file, original_filename=None):
    """解析上传的文件内容，兼容Streamlit UploadedFile和标准文件句柄"""
    content = ""
    filename_for_error = ""
    tmp_path = None  # 初始化tmp_path

    try:
        # 根据文件对象的类型获取文件名和内容
        if hasattr(file, 'getvalue'):  # Streamlit UploadedFile
            filename_for_type = file.name
            filename_for_error = file.name
            file_content = file.getvalue()
        else:  # 标准文件句柄 (from open(..., 'rb'))
            # 对于文件句柄，必须依赖传递的original_filename来获取类型
            filename_for_type = original_filename if original_filename else file.name
            filename_for_error = file.name # 在错误日志中记录实际路径
            file_content = file.read()

        # 使用os.path.splitext安全地获取文件扩展名
        _ , file_ext = os.path.splitext(filename_for_type)
        file_type = file_ext.lower().replace('.', '')

        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name

        if file_type == "txt":
            # 尝试用utf-8解码，如果失败则尝试其他编码
            try:
                with open(tmp_path, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                with open(tmp_path, "r", encoding="gbk", errors='ignore') as f:
                    content = f.read()

        elif file_type == "pdf":
            reader = PyPDF2.PdfReader(tmp_path)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    content += text + "\n"

        elif file_type == "docx":
            doc = Document(tmp_path)
            for para in doc.paragraphs:
                content += para.text + "\n"

        elif file_type == "pptx":
            prs = Presentation(tmp_path)  # 使用正确的类名
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"

        return content if content and content.strip() else ""

    except Exception as e:
        st.error(f"解析文件 '{filename_for_error}' 时出错: {str(e)}")
        return None
    finally:
        # 确保临时文件被删除
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)
    
def preprocess_text(text):
    """文本预处理"""
    if not text:
        return ""
    
    # 转换为小写
    text = text.lower()
    # 移除特殊字符和数字
    text = re.sub(r'[^a-zA-Z\u4e00-\u9fff\s]', '', text)
    # 分词
    tokens = word_tokenize(text)
    # 移除停用词
    stop_words = set(stopwords.words('english') + stopwords.words('chinese'))
    tokens = [word for word in tokens if word not in stop_words]
    return " ".join(tokens)

def split_into_chunks(text, chunk_size=200):
    """将文本分割为块"""
    if not text:
        return []
    
    chunks = []
    words = text.split()

    for i in range(0, len(words), chunk_size):
        chunks.append(" ".join(words[i:i + chunk_size]))

    return chunks

def generate_file_id(file_content):
    """生成文件唯一ID"""
    return hashlib.md5(file_content).hexdigest()